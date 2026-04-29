"""
模块 2: PPT 报告生成模块 (中英双语)
- 基于 python-pptx 生成 16:9 宽屏报告
- 支持项目名/规格信息头、斑马纹数据行、图片嵌入
- 自动分页（每页最多 7 行数据）
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class PPTGenerator:
    """VNA 测试报告 PPT 生成器"""

    def __init__(self, output_path):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

    def generate(self, dataset, proj_name="", spec="", lang="en"):
        """生成完整 PPT 报告"""
        blank_layout = self.prs.slide_layouts[6]

        # 颜色常量
        dark_blue = RGBColor(68, 114, 196)
        row_bg_1 = RGBColor(233, 237, 244)
        row_bg_2 = RGBColor(217, 225, 242)
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        gray = RGBColor(127, 140, 141)

        lbl_proj = "Project:" if lang == "en" else "项目名:"
        lbl_spec = "Spec:" if lang == "en" else "规格:"

        def format_cell(cell, text, bg_color, font_color, is_bold=False):
            """格式化单元格：背景色、字体、对齐方式"""
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.bold = is_bold
                paragraph.font.color.rgb = font_color
                paragraph.font.size = Pt(12)

        for sample_name, df in dataset.items():
            if df.empty:
                continue

            if lang == "en":
                sample_name = sample_name.replace("样品", "Sample ")

            chunk_size = 7
            chunks = [
                df[i : i + chunk_size] for i in range(0, len(df), chunk_size)
            ]

            for chunk_idx, chunk_df in enumerate(chunks):
                slide = self.prs.slides.add_slide(blank_layout)

                # --- 信息头区域 ---
                info_top = Inches(0.15)
                newlines = max(
                    proj_name.count("\n"), spec.count("\n")
                )
                info_height = Inches(0.3) + newlines * Inches(0.2)

                if proj_name or spec:
                    title_top = info_top + info_height + Inches(0.05)
                    info_box = slide.shapes.add_textbox(
                        Inches(0.5), info_top, Inches(15), info_height
                    )
                    tf = info_box.text_frame
                    tf.word_wrap = True
                    tf.text = f"{lbl_proj} {proj_name}      {lbl_spec} {spec}"
                    for p_info in tf.paragraphs:
                        p_info.font.size = Pt(13)
                        p_info.font.bold = True
                        p_info.font.color.rgb = gray
                else:
                    title_top = Inches(0.2)

                # --- 标题 ---
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), title_top, Inches(15), Inches(0.5)
                )
                p = title_box.text_frame.paragraphs[0]
                lbl_suffix = (
                    f" (Cont.{chunk_idx})"
                    if lang == "en"
                    else f" (续{chunk_idx})"
                )
                suffix = lbl_suffix if chunk_idx > 0 else ""
                p.text = sample_name + suffix
                p.font.size = Pt(24)
                p.font.bold = True

                # --- 表格 ---
                table_top = title_top + Inches(0.6)
                rows = len(chunk_df) + 1
                cols = 9
                header_height = Pt(45)

                max_bottom_margin = Inches(8.5)
                max_table_height = max_bottom_margin - table_top

                if len(chunk_df) <= 4:
                    data_row_height = Inches(1.6)
                    total_height = header_height + (
                        data_row_height * len(chunk_df)
                    )
                else:
                    available_data_height = max_table_height - header_height
                    data_row_height = int(
                        available_data_height / len(chunk_df)
                    )
                    total_height = max_table_height

                left, width = Inches(0.5), Inches(15)
                table_shape = slide.shapes.add_table(
                    rows, cols, left, table_top, width, total_height
                )
                table = table_shape.table

                table.rows[0].height = header_height
                for i in range(1, rows):
                    table.rows[i].height = int(data_row_height)

                # 列宽分配
                table.columns[0].width = Inches(1.0)
                for i in [1, 2, 3, 5, 6, 7]:
                    table.columns[i].width = Inches(1.3)
                for i in [4, 8]:
                    table.columns[i].width = Inches(3.1)

                # 表头
                headers = [
                    "Items",
                    "1.500GHz(IL)",
                    "3.000GHz(IL)",
                    "4.500GHz(IL)",
                    "Image",
                    "1.500GHz(RL)",
                    "3.000GHz(RL)",
                    "4.500GHz(RL)",
                    "Image",
                ]
                for col_idx, header in enumerate(headers):
                    format_cell(
                        table.cell(0, col_idx),
                        header,
                        dark_blue,
                        white,
                        is_bold=True,
                    )

                # 数据行
                for idx, (index, row_data) in enumerate(chunk_df.iterrows()):
                    row_idx = idx + 1
                    row_bg = row_bg_1 if row_idx % 2 == 1 else row_bg_2

                    pt_name = row_data["PointName"]
                    if lang == "en":
                        pt_name = pt_name.replace("点位", "Point ")

                    format_cell(
                        table.cell(row_idx, 0),
                        pt_name,
                        row_bg,
                        black,
                        is_bold=True,
                    )
                    format_cell(
                        table.cell(row_idx, 1),
                        row_data["1.5G_IL"],
                        row_bg,
                        black,
                    )
                    format_cell(
                        table.cell(row_idx, 2),
                        row_data["3.0G_IL"],
                        row_bg,
                        black,
                    )
                    format_cell(
                        table.cell(row_idx, 3),
                        row_data["4.5G_IL"],
                        row_bg,
                        black,
                    )
                    format_cell(table.cell(row_idx, 4), "", row_bg, black)
                    self._insert_image_to_cell(
                        slide, table_shape, row_idx, 4, row_data["Img_IL"]
                    )

                    format_cell(
                        table.cell(row_idx, 5),
                        row_data["1.5G_RL"],
                        row_bg,
                        black,
                    )
                    format_cell(
                        table.cell(row_idx, 6),
                        row_data["3.0G_RL"],
                        row_bg,
                        black,
                    )
                    format_cell(
                        table.cell(row_idx, 7),
                        row_data["4.5G_RL"],
                        row_bg,
                        black,
                    )
                    format_cell(table.cell(row_idx, 8), "", row_bg, black)
                    self._insert_image_to_cell(
                        slide, table_shape, row_idx, 8, row_data["Img_RL"]
                    )

        self.prs.save(self.output_path)

    def _insert_image_to_cell(self, slide, table_shape, row_idx, col_idx, img_path):
        """【核心算法】计算单元格绝对坐标并等比例插入图片"""
        if not img_path or not os.path.exists(img_path):
            return

        table = table_shape.table
        cell_x = table_shape.left + sum(
            table.columns[i].width for i in range(col_idx)
        )
        cell_y = table_shape.top + sum(
            table.rows[j].height for j in range(row_idx)
        )
        cell_w = table.columns[col_idx].width
        cell_h = table.rows[row_idx].height

        try:
            from PIL import Image

            with Image.open(img_path) as img:
                img_w, img_h = img.size
        except Exception:
            return

        ratio = min(cell_w / img_w * 0.95, cell_h / img_h * 0.90)
        fit_w, fit_h = int(img_w * ratio), int(img_h * ratio)
        offset_x = cell_x + (cell_w - fit_w) / 2
        offset_y = cell_y + (cell_h - fit_h) / 2
        slide.shapes.add_picture(img_path, offset_x, offset_y, fit_w, fit_h)
