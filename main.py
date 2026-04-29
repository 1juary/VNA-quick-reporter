import sys
import os
import re
import json
import cv2
import pandas as pd
from pathlib import Path
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import numpy as np

from rapidocr_onnxruntime import RapidOCR

from PySide6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, 
    QPushButton, QLabel, QFileDialog, QTabWidget, QTableWidget, 
    QHeaderView, QProgressDialog, QMessageBox, QDialog, QTextEdit,
    QLineEdit, QComboBox, QFrame
)
from PySide6.QtCore import Qt, QThread, Signal, QMimeData
from PySide6.QtGui import QPixmap, QDrag

CONFIG_FILE = "vna_config.json"

# ==========================================
# 模块 1: RapidOCR 图像处理与提取模块
# ==========================================
class VNAOCRExtractor:
    def __init__(self):
        self.ocr = RapidOCR()
        self.target_freqs =[1.5, 3.0, 4.5]
        self.pattern = re.compile(r'(\d+\.\d+)\s*[Gg][Hh][Zz][^\d-]{0,25}(-?\d+\.\d+)\s*[dD][Bb]')

    def process_image(self, img_path):
        if not img_path or not os.path.exists(img_path):
            return {1.5: "N/A", 3.0: "N/A", 4.5: "N/A"}

        try:
            img = cv2.imdecode(np.fromfile(img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
        except Exception:
            img = None

        if img is None:
            return {1.5: "N/A", 3.0: "N/A", 4.5: "N/A"}

        h, w = img.shape[:2]
        
        crop_default = img[int(h * 0.10):int(h * 0.28), int(w * 0.65):int(w * 0.88)]
        res_default = self._extract(crop_default)
        if self._is_complete(res_default):
            return res_default
            
        crop_shifted = img[int(h * 0.05):int(h * 0.35), int(w * 0.35):int(w * 0.80)]
        res_shifted = self._extract(crop_shifted)
        if self._has_any_data(res_shifted):
            return res_shifted
            
        crop_extreme = img[0:int(h * 0.50), 0:w]
        return self._extract(crop_extreme)

    def _extract(self, crop):
        result, _ = self.ocr(crop)
        text_content = ""
        
        if result:
            for line in result:
                text_content += str(line[1]) + " "
                
        results = {1.5: "N/A", 3.0: "N/A", 4.5: "N/A"}
        matches = self.pattern.findall(text_content)
        
        for match in matches:
            try:
                freq = float(match[0])
                db_val = f"{float(match[1]):.2f}dB" 
                closest_freq = min(self.target_freqs, key=lambda x: abs(x - freq))
                if abs(closest_freq - freq) < 0.2:
                    if results[closest_freq] == "N/A":
                        results[closest_freq] = db_val
            except ValueError:
                continue
        return results

    def _is_complete(self, res):
        return sum(1 for v in res.values() if v != "N/A") == len(self.target_freqs)
        
    def _has_any_data(self, res):
        return sum(1 for v in res.values() if v != "N/A") > 0

# ==========================================
# 模块 2: PPT 报告生成模块 (中英双语)
# ==========================================
class PPTGenerator:
    def __init__(self, output_path):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

    def generate(self, dataset, proj_name="", spec="", lang="en"):
        blank_layout = self.prs.slide_layouts[6] 

        dark_blue = RGBColor(68, 114, 196)   
        row_bg_1 = RGBColor(233, 237, 244)   
        row_bg_2 = RGBColor(217, 225, 242)   
        white = RGBColor(255, 255, 255)
        black = RGBColor(0, 0, 0)
        gray = RGBColor(127, 140, 141)

        lbl_proj = "Project:" if lang == "en" else "项目名:"
        lbl_spec = "Spec:" if lang == "en" else "规格:"

        def format_cell(cell, text, bg_color, font_color, is_bold=False):
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE 
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER 
                paragraph.font.bold = is_bold
                paragraph.font.color.rgb = font_color
                paragraph.font.size = Pt(12) 

        for sample_name, df in dataset.items():
            if df.empty: continue
            
            if lang == "en":
                sample_name = sample_name.replace("样品", "Sample ")
            
            chunk_size = 7
            chunks = [df[i:i + chunk_size] for i in range(0, len(df), chunk_size)]
            
            for chunk_idx, chunk_df in enumerate(chunks):
                slide = self.prs.slides.add_slide(blank_layout)
                
                info_top = Inches(0.15)
                newlines = max(proj_name.count('\n'), spec.count('\n'))
                info_height = Inches(0.3) + newlines * Inches(0.2)
                
                if proj_name or spec:
                    title_top = info_top + info_height + Inches(0.05)
                    info_box = slide.shapes.add_textbox(Inches(0.5), info_top, Inches(15), info_height)
                    tf = info_box.text_frame
                    tf.word_wrap = True 
                    tf.text = f"{lbl_proj} {proj_name}      {lbl_spec} {spec}"
                    
                    for p_info in tf.paragraphs:
                        p_info.font.size = Pt(13)
                        p_info.font.bold = True
                        p_info.font.color.rgb = gray
                else:
                    title_top = Inches(0.2)

                title_box = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(15), Inches(0.5))
                p = title_box.text_frame.paragraphs[0]
                
                lbl_suffix = f" (Cont.{chunk_idx})" if lang == "en" else f" (续{chunk_idx})"
                suffix = lbl_suffix if chunk_idx > 0 else ""
                
                p.text = sample_name + suffix
                p.font.size = Pt(24)
                p.font.bold = True

                table_top = title_top + Inches(0.6)
                rows = len(chunk_df) + 1 
                cols = 9
                header_height = Pt(45) 
                
                max_bottom_margin = Inches(8.5) 
                max_table_height = max_bottom_margin - table_top
                
                if len(chunk_df) <= 4:
                    data_row_height = Inches(1.6) 
                    total_height = header_height + (data_row_height * len(chunk_df))
                else:
                    available_data_height = max_table_height - header_height
                    data_row_height = int(available_data_height / len(chunk_df))
                    total_height = max_table_height

                left, width = Inches(0.5), Inches(15)
                table_shape = slide.shapes.add_table(rows, cols, left, table_top, width, total_height)
                table = table_shape.table
                
                table.rows[0].height = header_height
                for i in range(1, rows): 
                    table.rows[i].height = int(data_row_height)

                table.columns[0].width = Inches(1.0)
                for i in[1, 2, 3, 5, 6, 7]: table.columns[i].width = Inches(1.3)
                for i in [4, 8]: table.columns[i].width = Inches(3.1)

                headers =["Items", "1.500GHz(IL)", "3.000GHz(IL)", "4.500GHz(IL)", "Image", 
                           "1.500GHz(RL)", "3.000GHz(RL)", "4.500GHz(RL)", "Image"]
                for col_idx, header in enumerate(headers):
                    format_cell(table.cell(0, col_idx), header, dark_blue, white, is_bold=True)

                for idx, (index, row_data) in enumerate(chunk_df.iterrows()):
                    row_idx = idx + 1 
                    row_bg = row_bg_1 if row_idx % 2 == 1 else row_bg_2
                    
                    pt_name = row_data['PointName']
                    if lang == "en":
                        pt_name = pt_name.replace("点位", "Point ")
                        
                    format_cell(table.cell(row_idx, 0), pt_name, row_bg, black, is_bold=True)
                    
                    format_cell(table.cell(row_idx, 1), row_data['1.5G_IL'], row_bg, black)
                    format_cell(table.cell(row_idx, 2), row_data['3.0G_IL'], row_bg, black)
                    format_cell(table.cell(row_idx, 3), row_data['4.5G_IL'], row_bg, black)
                    format_cell(table.cell(row_idx, 4), "", row_bg, black) 
                    self._insert_image_to_cell(slide, table_shape, row_idx, 4, row_data['Img_IL'])
                    
                    format_cell(table.cell(row_idx, 5), row_data['1.5G_RL'], row_bg, black)
                    format_cell(table.cell(row_idx, 6), row_data['3.0G_RL'], row_bg, black)
                    format_cell(table.cell(row_idx, 7), row_data['4.5G_RL'], row_bg, black)
                    format_cell(table.cell(row_idx, 8), "", row_bg, black) 
                    self._insert_image_to_cell(slide, table_shape, row_idx, 8, row_data['Img_RL'])

        self.prs.save(self.output_path)

    def _insert_image_to_cell(self, slide, table_shape, row_idx, col_idx, img_path):
        if not img_path or not os.path.exists(img_path): return
        table = table_shape.table
        cell_x = table_shape.left + sum(table.columns[i].width for i in range(col_idx))
        cell_y = table_shape.top + sum(table.rows[j].height for j in range(row_idx))
        cell_w = table.columns[col_idx].width
        cell_h = table.rows[row_idx].height

        try:
            from PIL import Image
            with Image.open(img_path) as img:
                img_w, img_h = img.size
        except Exception: return

        ratio = min(cell_w / img_w * 0.95, cell_h / img_h * 0.90)
        fit_w, fit_h = int(img_w * ratio), int(img_h * ratio)
        offset_x = cell_x + (cell_w - fit_w) / 2
        offset_y = cell_y + (cell_h - fit_h) / 2
        slide.shapes.add_picture(img_path, offset_x, offset_y, fit_w, fit_h)

# ==========================================
# 模块 3: 自定义 UI 组件 (支持相互拖动互换 + 一键删除)
# ==========================================
class ImageCell(QLabel):
    imageLoaded = Signal(str)
    filesDroppedToTab = Signal(list) 

    def __init__(self):
        super().__init__()
        self.image_path = ""
        self.setAcceptDrops(True)
        self.drag_start_pos = None
        
        self.btn_layout = QVBoxLayout(self)
        self.btn_layout.setContentsMargins(5, 5, 5, 5)
        # 【修改点】修正为 Qt.AlignTop | Qt.AlignRight
        self.btn_layout.setAlignment(Qt.AlignTop | Qt.AlignRight)
        
        self.btn_delete = QPushButton("×")
        self.btn_delete.setFixedSize(22, 22)
        self.btn_delete.setCursor(Qt.PointingHandCursor)
        self.btn_delete.setStyleSheet("""
            QPushButton {
                background-color: #E74C3C;
                color: white;
                border-radius: 11px;
                font-weight: bold;
                font-family: Arial;
                font-size: 14px;
                padding-bottom: 2px;
            }
            QPushButton:hover { background-color: #C0392B; }
        """)
        self.btn_delete.clicked.connect(self.clear_image)
        self.btn_delete.hide()
        
        self.btn_layout.addWidget(self.btn_delete)
        self.reset_ui()

    def reset_ui(self):
        self.setText("点击 或 拖拽\n(按住图片可互相交换)")
        self.setAlignment(Qt.AlignCenter)
        self.setCursor(Qt.PointingHandCursor) 
        self.setStyleSheet("""
            QLabel {
                background-color: #F8F9FA; color: #95A5A6;
                font-size: 13px; font-weight: bold;
                border: 2px dashed #D1D8E0; border-radius: 6px; margin: 4px;
            }
            QLabel:hover { background-color: #E2E8F0; border-color: #3DC2EC; color: #2C3E50; }
        """)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.drag_start_pos = event.position().toPoint()

    def mouseMoveEvent(self, event):
        if not self.drag_start_pos or not self.image_path:
            return
        if (event.position().toPoint() - self.drag_start_pos).manhattanLength() > QApplication.startDragDistance():
            drag = QDrag(self)
            mime = QMimeData()
            mime.setText(self.image_path)
            drag.setMimeData(mime)
            drag.setPixmap(self.pixmap().scaled(100, 100, Qt.KeepAspectRatio))
            drag.exec(Qt.MoveAction)
            self.drag_start_pos = None

    def mouseReleaseEvent(self, event):
        if self.drag_start_pos is not None:
            path, _ = QFileDialog.getOpenFileName(self, "选择图片", "", "Images (*.png *.jpg *.jpeg)")
            if path: self.load_image(path)
            self.drag_start_pos = None

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls() or event.mimeData().hasText():
            event.acceptProposedAction()

    def dropEvent(self, event):
        source = event.source()
        if isinstance(source, ImageCell) and source != self:
            source_path = source.image_path
            target_path = self.image_path
            
            if source_path: self.load_image(source_path)
            else: self.clear_image()
            
            if target_path: source.load_image(target_path)
            else: source.clear_image()
            
            event.acceptProposedAction()
            
        elif event.mimeData().hasUrls():
            urls = event.mimeData().urls()
            if len(urls) == 1 and Path(urls[0].toLocalFile()).is_file():
                self.load_image(urls[0].toLocalFile())
                event.acceptProposedAction()
            else:
                self.filesDroppedToTab.emit(urls)
                event.acceptProposedAction()

    def load_image(self, path):
        self.image_path = str(path)
        pixmap = QPixmap(self.image_path)
        self.setPixmap(pixmap.scaled(self.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation))
        self.setStyleSheet("border: none; background-color: transparent; margin: 4px;")
        self.btn_delete.show() 
        self.imageLoaded.emit(self.image_path)

    def clear_image(self):
        self.image_path = ""
        self.setPixmap(QPixmap())
        self.btn_delete.hide() 
        self.reset_ui()

# ==========================================
# 模块 4: GUI 业务流线程与界面
# ==========================================
class OCRWorker(QThread):
    progress_update = Signal(int, str)
    finished = Signal(dict)

    def __init__(self, ui_data):
        super().__init__()
        self.ui_data = ui_data 
        
    def run(self):
        extractor = VNAOCRExtractor()
        result_dataset = {}
        total_tasks = sum(len(pairs) for pairs in self.ui_data.values())
        current_task = 0

        for sample_name, pairs in self.ui_data.items():
            sample_data =[]
            for pair in pairs:
                current_task += 1
                self.progress_update.emit(int(current_task / total_tasks * 100), f"正在处理: {sample_name}...")
                
                il_data = extractor.process_image(pair['IL'])
                rl_data = extractor.process_image(pair['RL'])
                
                sample_data.append({
                    'PointName': pair['PointName'], 
                    '1.5G_IL': il_data.get(1.5),
                    '3.0G_IL': il_data.get(3.0),
                    '4.5G_IL': il_data.get(4.5),
                    'Img_IL': pair['IL'],
                    '1.5G_RL': rl_data.get(1.5),
                    '3.0G_RL': rl_data.get(3.0),
                    '4.5G_RL': rl_data.get(4.5),
                    'Img_RL': pair['RL']
                })
            result_dataset[sample_name] = pd.DataFrame(sample_data)
        self.finished.emit(result_dataset)

class SampleTab(QWidget):
    def __init__(self):
        super().__init__()
        self.layout = QVBoxLayout(self)
        self.layout.setContentsMargins(15, 15, 15, 15)
        
        self.table = QTableWidget(3, 3)
        self.table.setHorizontalHeaderLabels(["测试点位", "IL IMAGE", "RL IMAGE"])
        self.table.horizontalHeader().setSectionResizeMode(0, QHeaderView.Fixed)
        self.table.horizontalHeader().setSectionResizeMode(1, QHeaderView.Stretch)
        self.table.horizontalHeader().setSectionResizeMode(2, QHeaderView.Stretch)
        self.table.setColumnWidth(0, 110) 
        self.table.verticalHeader().setDefaultSectionSize(120)
        self.table.setSelectionMode(QTableWidget.NoSelection)
        self.table.setFocusPolicy(Qt.NoFocus)

        for row in range(3): self.init_row_widgets(row)
        self.layout.addWidget(self.table)
        
        self.btn_add_row = QPushButton("+ 添加点位")
        self.btn_add_row.setObjectName("BtnAddRow")
        self.btn_add_row.clicked.connect(self.add_row)
        self.layout.addWidget(self.btn_add_row, alignment=Qt.AlignCenter)

    def init_row_widgets(self, row):
        point_edit = QLineEdit(f"点位{row+1}")
        point_edit.setAlignment(Qt.AlignCenter)
        point_edit.setStyleSheet("border: none; background: transparent; font-weight: bold; color: #2C3E50; font-size: 14px;")
        self.table.setCellWidget(row, 0, point_edit)
        
        il_cell = ImageCell()
        il_cell.imageLoaded.connect(lambda path, r=row: self.auto_fill_point_name(r, path))
        il_cell.filesDroppedToTab.connect(self.handle_dropped_files) 
        self.table.setCellWidget(row, 1, il_cell)
        
        rl_cell = ImageCell()
        rl_cell.imageLoaded.connect(lambda path, r=row: self.auto_fill_point_name(r, path))
        rl_cell.filesDroppedToTab.connect(self.handle_dropped_files)
        self.table.setCellWidget(row, 2, rl_cell)

    def add_row(self):
        row = self.table.rowCount()
        self.table.insertRow(row)
        self.init_row_widgets(row)

    def auto_fill_point_name(self, row, path):
        stem = Path(path).stem
        clean_name = re.sub(r'[-_]([iI][lL]|[rR][lL]|\d+)$', '', stem)
        edit_widget = self.table.cellWidget(row, 0)
        current_text = edit_widget.text().strip()
        if current_text.startswith("点位") or not current_text:
            edit_widget.setText(clean_name)

    def handle_dropped_files(self, urls):
        files =[]
        for url in urls:
            path = Path(url.toLocalFile())
            if path.is_dir(): files.extend([str(p) for p in path.rglob("*") if p.is_file()])
            elif path.is_file(): files.append(str(path))
                
        img_files = sorted([f for f in files if f.lower().endswith(('.jpg', '.jpeg', '.png'))])
        if not img_files: return

        groups = defaultdict(list)
        for f in img_files:
            stem = Path(f).stem
            base_name = re.sub(r'[-_]([iI][lL]|[rR][lL]|\d+)$', '', stem)
            groups[base_name].append(f)

        final_pairs = []
        orphans =[]

        for base, f_list in groups.items():
            f_list = sorted(f_list)
            
            def intelligent_sort(x):
                xl = x.lower()
                if 'il' in xl and 'rl' not in xl: return 0
                if 'rl' in xl and 'il' not in xl: return 2
                return 1
            f_list.sort(key=intelligent_sort)

            while len(f_list) >= 2:
                final_pairs.append((base, f_list[0], f_list[1]))
                f_list = f_list[2:]
            
            if f_list: orphans.extend(f_list)

        orphans = sorted(orphans)
        while len(orphans) >= 2:
            o1, o2 = orphans[0], orphans[1]
            stem = Path(o1).stem
            base = re.sub(r'[-_]([iI][lL]|[rR][lL]|\d+)$', '', stem)
            final_pairs.append((base, o1, o2))
            orphans = orphans[2:]

        for base, il_path, rl_path in final_pairs:
            row = self.find_empty_row_or_add()
            self.table.cellWidget(row, 0).setText(base)
            self.table.cellWidget(row, 1).load_image(il_path)
            self.table.cellWidget(row, 2).load_image(rl_path)
            
        if orphans:
            QMessageBox.information(self, "警告", f"由于文件数为奇数，有 {len(orphans)} 个文件未能配对！")

    def find_empty_row_or_add(self):
        for row in range(self.table.rowCount()):
            if not self.table.cellWidget(row, 1).image_path and not self.table.cellWidget(row, 2).image_path:
                return row
        row = self.table.rowCount()
        self.add_row()
        return row

    def get_image_pairs(self):
        pairs =[]
        for row in range(self.table.rowCount()):
            point_name = self.table.cellWidget(row, 0).text().strip()
            il_cell = self.table.cellWidget(row, 1)
            rl_cell = self.table.cellWidget(row, 2)
            if il_cell.image_path and rl_cell.image_path:
                pairs.append({
                    'PointName': point_name if point_name else f"点位{row+1}",
                    'IL': il_cell.image_path,
                    'RL': rl_cell.image_path
                })
        return pairs

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("VNA Data Automator")
        self.resize(1000, 780)
        self.setAcceptDrops(True) 
        self.config_map = self.load_settings()
        self.init_ui()
        self.apply_styles()

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls(): event.accept()
        else: event.ignore()
            
    def dropEvent(self, event):
        urls = event.mimeData().urls()
        current_tab = self.tabs.currentWidget()
        if hasattr(current_tab, 'handle_dropped_files'):
            current_tab.handle_dropped_files(urls)

    def load_settings(self):
        if os.path.exists(CONFIG_FILE):
            try:
                with open(CONFIG_FILE, 'r', encoding='utf-8') as f: return json.load(f)
            except Exception: pass
        return {}

    def save_settings(self, proj, spec):
        if proj:
            self.config_map[proj] = spec
            try:
                with open(CONFIG_FILE, 'w', encoding='utf-8') as f:
                    json.dump(self.config_map, f, ensure_ascii=False, indent=2)
            except Exception: pass

    def init_ui(self):
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        main_layout = QVBoxLayout(main_widget)
        main_layout.setSpacing(15)
        
        top_card = QFrame()
        top_card.setObjectName("InfoCard")
        top_layout = QHBoxLayout(top_card)
        top_layout.setContentsMargins(20, 15, 20, 15)
        
        lbl_lang = QLabel("报告语言 :")
        lbl_lang.setObjectName("InfoTitle")
        self.combo_lang = QComboBox()
        self.combo_lang.addItems(["English", "中文"])
        self.combo_lang.setFocusPolicy(Qt.NoFocus)

        lbl_proj = QLabel("项目名 :")
        lbl_proj.setObjectName("InfoTitle")
        self.edit_proj = QLineEdit()
        self.edit_proj.setPlaceholderText("例如: Project X")
        self.edit_proj.textChanged.connect(self.on_project_name_changed)
        
        lbl_spec = QLabel("规 格 :")
        lbl_spec.setObjectName("InfoTitle")
        self.combo_spec = QComboBox()
        self.combo_spec.setEditable(True)
        self.combo_spec.setPlaceholderText("请选择或手动输入规格")
        self.combo_spec.addItems(list(set(self.config_map.values())))

        top_layout.addWidget(lbl_lang)
        top_layout.addWidget(self.combo_lang, stretch=1)
        top_layout.addSpacing(20)
        top_layout.addWidget(lbl_proj)
        top_layout.addWidget(self.edit_proj, stretch=2)
        top_layout.addSpacing(20)
        top_layout.addWidget(lbl_spec)
        top_layout.addWidget(self.combo_spec, stretch=2)
        main_layout.addWidget(top_card)

        self.tabs = QTabWidget()
        self.tabs.addTab(SampleTab(), "样品1")
        self.tabs.addTab(SampleTab(), "样品2")
        self.tabs.addTab(SampleTab(), "样品3")
        
        self.btn_add_tab = QPushButton("+ 新增样品")
        self.btn_add_tab.setObjectName("ChartButton") 
        self.btn_add_tab.setCursor(Qt.PointingHandCursor)
        self.btn_add_tab.clicked.connect(lambda: self.tabs.addTab(SampleTab(), f"样品{self.tabs.count()+1}"))
        self.tabs.setCornerWidget(self.btn_add_tab, Qt.TopRightCorner)
        main_layout.addWidget(self.tabs)

        bottom_layout = QHBoxLayout()
        bottom_layout.addStretch()
        
        self.btn_preview = QPushButton("预览数据")
        self.btn_preview.setObjectName("ActionButton")
        self.btn_preview.setCursor(Qt.PointingHandCursor)
        self.btn_preview.clicked.connect(self.preview_data)
        
        self.btn_export = QPushButton("🚀 导出完整报告")
        self.btn_export.setObjectName("ActionButton")
        self.btn_export.setCursor(Qt.PointingHandCursor)
        self.btn_export.clicked.connect(self.export_ppt)
        
        bottom_layout.addWidget(self.btn_preview)
        bottom_layout.addSpacing(15)
        bottom_layout.addWidget(self.btn_export)
        bottom_layout.addStretch()
        main_layout.addLayout(bottom_layout)

    def on_project_name_changed(self, text):
        text = text.strip()
        if text in self.config_map:
            self.combo_spec.setCurrentText(self.config_map[text])

    def apply_styles(self):
        self.setStyleSheet("""
            QMainWindow { background-color: #F4F6F8; }
            QFrame#InfoCard { background-color: #FFFFFF; border-radius: 12px; border: 1px solid #D1D8E0; }
            QLabel#InfoTitle { font-size: 14px; font-weight: 800; color: #2C3E50; }
            QLineEdit, QComboBox { border: 1px solid #D1D8E0; border-radius: 6px; padding: 6px 10px; font-size: 13px; color: #2C3E50; background-color: #F8F9FA; }
            QLineEdit:focus, QComboBox:focus { border-color: #3DC2EC; background-color: #FFFFFF;}
            QTabWidget::pane { border: 1px solid #D1D8E0; border-radius: 8px; background-color: #FFFFFF; }
            QTabBar::tab { background-color: #E2E8F0; color: #7F8C8D; padding: 10px 25px; margin-right: 4px; border-top-left-radius: 8px; border-top-right-radius: 8px; font-weight: bold; font-size: 14px; }
            QTabBar::tab:selected { background-color: #FFFFFF; color: #2C3E50; border: 1px solid #D1D8E0; border-bottom: none; }
            QTableWidget { background-color: #FFFFFF; border: none; gridline-color: #E9ECEF; }
            QHeaderView::section { background-color: #F8F9FA; color: #34495E; font-weight: bold; font-size: 14px; border: none; border-bottom: 2px solid #D1D8E0; padding: 12px; }
            QPushButton#BtnAddRow { background-color: transparent; font-size: 14px; font-weight: bold; color: #95A5A6; padding: 10px; }
            QPushButton#BtnAddRow:hover { color: #3DC2EC; }
            QPushButton#ChartButton { background-color: transparent; border-radius: 8px; font-size: 13px; font-weight: bold; color: #34495E; padding: 5px 15px; margin-top: 5px; }
            QPushButton#ChartButton:hover { background-color: #E2E8F0; }
            QPushButton#ActionButton { background-color: #3DC2EC; color: #FFFFFF; font-size: 15px; font-weight: bold; border: none; border-radius: 20px; padding: 12px 35px; }
            QPushButton#ActionButton:hover { background-color: #5ED1F4; }
            QPushButton#ActionButton:pressed { background-color: #2BAAD4; }
        """)

    def gather_ui_data(self):
        ui_data = {}
        for i in range(self.tabs.count()):
            tab_name = self.tabs.tabText(i)
            tab_widget = self.tabs.widget(i)
            pairs = tab_widget.get_image_pairs()
            if pairs: ui_data[tab_name] = pairs
        return ui_data

    def preview_data(self):
        ui_data = self.gather_ui_data()
        if not ui_data:
            QMessageBox.warning(self, "提示", "请至少放入一组完整的 IL / RL 图片！")
            return
        self.start_ocr_task(ui_data, mode="preview")

    def export_ppt(self):
        ui_data = self.gather_ui_data()
        if not ui_data:
            QMessageBox.warning(self, "提示", "请至少放入一组完整的 IL / RL 图片！")
            return
            
        save_path, _ = QFileDialog.getSaveFileName(self, "保存 PPT", "VNA_Report.pptx", "PowerPoint (*.pptx)")
        if not save_path: return

        proj = self.edit_proj.text().strip()
        spec = self.combo_spec.currentText().strip()
        self.save_settings(proj, spec)

        self.save_path = save_path
        self.start_ocr_task(ui_data, mode="export")

    def start_ocr_task(self, ui_data, mode):
        self.mode = mode
        self.current_lang = "en" if self.combo_lang.currentText() == "English" else "zh"
        
        self.progress_dialog = QProgressDialog("正在通过 RapidOCR 极速提取数据...", "取消", 0, 100, self)
        self.progress_dialog.setWindowTitle("处理中")
        self.progress_dialog.setWindowModality(Qt.WindowModal)
        self.progress_dialog.show()

        self.worker = OCRWorker(ui_data)
        self.worker.progress_update.connect(self.progress_dialog.setValue)
        self.worker.progress_update.connect(lambda v, t: self.progress_dialog.setLabelText(t))
        self.worker.finished.connect(self.on_ocr_finished)
        self.worker.start()

    def on_ocr_finished(self, result_dataset):
        self.progress_dialog.setValue(100)
        
        if self.mode == "preview":
            self.show_preview_dialog(result_dataset)
        elif self.mode == "export":
            try:
                ppt_gen = PPTGenerator(self.save_path)
                ppt_gen.generate(
                    result_dataset, 
                    proj_name=self.edit_proj.text().strip(), 
                    spec=self.combo_spec.currentText().strip(),
                    lang=self.current_lang 
                )
                QMessageBox.information(self, "成功", f"PPT已成功导出至:\n{self.save_path}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"PPT生成失败: {str(e)}")

    def show_preview_dialog(self, result_dataset):
        dialog = QDialog(self)
        dialog.setWindowTitle("数据提取预览")
        dialog.resize(800, 400)
        layout = QVBoxLayout(dialog)
        
        preview_text = ""
        for sample, df in result_dataset.items():
            preview_text += f"=== {sample} ===\n"
            cols_to_show =[c for c in df.columns if 'Img' not in c]
            preview_text += df[cols_to_show].to_string() + "\n\n"
            
        label = QLabel("提取的数据如下（核对无误后请关闭此窗口并点击导出）：")
        text_edit = QTextEdit()
        text_edit.setPlainText(preview_text)
        text_edit.setReadOnly(True)
        text_edit.setStyleSheet("font-family: Consolas; font-size: 14px;")
        
        layout.addWidget(label)
        layout.addWidget(text_edit)
        dialog.exec()

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())